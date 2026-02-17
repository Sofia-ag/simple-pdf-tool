from flask import Flask, render_template, request, send_file
from pypdf import PdfReader, PdfWriter
from PIL import Image
from pdf2image import convert_from_path
import os
import uuid
import zipfile
import subprocess

# For watermark
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from io import BytesIO

# For converters
import pdfplumber
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "outputs"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# -------------------- PAGES --------------------

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/merge_page")
def merge_page():
    return render_template("merge.html")

@app.route("/split_page")
def split_page():
    return render_template("split.html")

@app.route("/compress_page")
def compress_page():
    return render_template("compress.html")

@app.route("/rotate_page")
def rotate_page():
    return render_template("rotate.html")

@app.route("/watermark_page")
def watermark_page():
    return render_template("watermark.html")

@app.route("/img_to_pdf_page")
def img_to_pdf_page():
    return render_template("img_to_pdf.html")

@app.route("/pdf_to_images_page")
def pdf_to_images_page():
    return render_template("pdf_to_images.html")

@app.route("/convert_page")
def convert_page():
    return render_template("convert.html")

# -------------------- ACTIONS --------------------

@app.route("/merge", methods=["POST"])
def merge_pdfs():
    files = request.files.getlist("pdfs")
    writer = PdfWriter()

    for file in files:
        path = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(path)
        reader = PdfReader(path)
        for page in reader.pages:
            writer.add_page(page)

    output_name = f"merged_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)

    with open(output_path, "wb") as f:
        writer.write(f)

    return send_file(output_path, as_attachment=True)

@app.route("/split", methods=["POST"])
def split_pdf():
    file = request.files["pdf"]
    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    reader = PdfReader(pdf_path)

    zip_name = f"split_{uuid.uuid4().hex}.zip"
    zip_path = os.path.join(OUTPUT_FOLDER, zip_name)

    with zipfile.ZipFile(zip_path, "w") as zipf:
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)

            out_name = f"page_{i+1}.pdf"
            out_path = os.path.join(OUTPUT_FOLDER, out_name)
            with open(out_path, "wb") as f:
                writer.write(f)

            zipf.write(out_path, out_name)
            os.remove(out_path)

    return send_file(zip_path, as_attachment=True)

@app.route("/rotate", methods=["POST"])
def rotate_pdf():
    file = request.files["pdf"]
    angle = int(request.form["angle"])

    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)

    output_name = f"rotated_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)

    with open(output_path, "wb") as f:
        writer.write(f)

    return send_file(output_path, as_attachment=True)

@app.route("/watermark", methods=["POST"])
def watermark_pdf():
    file = request.files["pdf"]
    text = request.form["text"]

    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    packet = BytesIO()
    can = canvas.Canvas(packet, pagesize=A4)
    can.setFont("Helvetica", 40)
    can.setFillAlpha(0.3)
    can.drawCentredString(300, 400, text)
    can.save()

    packet.seek(0)
    watermark_pdf = PdfReader(packet)

    for page in reader.pages:
        page.merge_page(watermark_pdf.pages[0])
        writer.add_page(page)

    output_name = f"watermarked_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)

    with open(output_path, "wb") as f:
        writer.write(f)

    return send_file(output_path, as_attachment=True)

# Basic python-only compress (stable, no Ghostscript)
@app.route("/compress", methods=["POST"])
def compress_pdf():
    file = request.files["pdf"]
    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    output_name = f"compressed_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)

    with open(output_path, "wb") as f:
        writer.write(f)

    return send_file(output_path, as_attachment=True)

@app.route("/img_to_pdf", methods=["POST"])
def img_to_pdf():
    files = request.files.getlist("images")
    images = []

    for file in files:
        img = Image.open(file.stream).convert("RGB")
        images.append(img)

    output_name = f"images_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)

    images[0].save(output_path, save_all=True, append_images=images[1:])
    return send_file(output_path, as_attachment=True)

@app.route("/pdf_to_images", methods=["POST"])
def pdf_to_images():
    file = request.files["pdf"]
    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    pages = convert_from_path(pdf_path)

    zip_name = f"images_{uuid.uuid4().hex}.zip"
    zip_path = os.path.join(OUTPUT_FOLDER, zip_name)

    with zipfile.ZipFile(zip_path, "w") as zipf:
        for i, page in enumerate(pages):
            img_name = f"page_{i+1}.png"
            img_path = os.path.join(OUTPUT_FOLDER, img_name)
            page.save(img_path, "PNG")
            zipf.write(img_path, img_name)
            os.remove(img_path)

    return send_file(zip_path, as_attachment=True)

# -------------------- CONVERTERS --------------------

# LibreOffice path (CHANGE if needed)
SOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"

def convert_to_pdf_with_libreoffice(input_path, output_dir):
    cmd = [
        SOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_path
    ]
    subprocess.run(cmd, check=True)

# Word / PPT / Excel / HTML -> PDF
@app.route("/office_to_pdf", methods=["POST"])
def office_to_pdf():
    file = request.files["file"]
    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(input_path)

    convert_to_pdf_with_libreoffice(input_path, OUTPUT_FOLDER)

    base = os.path.splitext(file.filename)[0]
    output_path = os.path.join(OUTPUT_FOLDER, base + ".pdf")
    return send_file(output_path, as_attachment=True)

# PDF -> Word (text-based)
@app.route("/pdf_to_word", methods=["POST"])
def pdf_to_word():
    file = request.files["pdf"]
    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    doc = Document()

    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            for line in text.split("\n"):
                doc.add_paragraph(line)

    output_name = f"pdf_to_word_{uuid.uuid4().hex}.docx"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)
    doc.save(output_path)

    return send_file(output_path, as_attachment=True)

# PDF -> PowerPoint (each page as image slide)
@app.route("/pdf_to_ppt", methods=["POST"])
def pdf_to_ppt():
    file = request.files["pdf"]
    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    pages = convert_from_path(pdf_path)
    prs = Presentation()

    for img in pages:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        img_name = f"tmp_{uuid.uuid4().hex}.png"
        img_path = os.path.join(OUTPUT_FOLDER, img_name)
        img.save(img_path, "PNG")

        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
        os.remove(img_path)

    output_name = f"pdf_to_ppt_{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)
    prs.save(output_path)

    return send_file(output_path, as_attachment=True)

# PDF -> Excel (text extraction)
@app.route("/pdf_to_excel", methods=["POST"])
def pdf_to_excel():
    file = request.files["pdf"]
    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    wb = Workbook()
    ws = wb.active
    ws.title = "Extracted"

    row_idx = 1

    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            for line in text.split("\n"):
                ws.cell(row=row_idx, column=1, value=line)
                row_idx += 1

    output_name = f"pdf_to_excel_{uuid.uuid4().hex}.xlsx"
    output_path = os.path.join(OUTPUT_FOLDER, output_name)
    wb.save(output_path)

    return send_file(output_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
